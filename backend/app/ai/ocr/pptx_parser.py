# =============================================================================
# File: pptx_parser.py
# Module/Service: Pipeline Worker — Document Parsing & Cleaning ([AI])
# Layer: Service
# Purpose: PPTX slide extraction into ordered _ParsedBlock units.
# Responsibilities:
#   - Extract title/subtitle/textbox/table/notes in visual order
#   - Map slide index to physical page_number
# Dependencies:
#   - python-pptx; app.ai.ocr.models/tables
# Public Exports:
#   - _parse_pptx, _parse_pptx_slide, _pptx_shape_blocks,
#     _pptx_table_to_text, _pptx_notes_text
# Database/Table: N/A
# Related Modules: app.ai.ocr.*
# Important Notes: page_count = slide count.
# =============================================================================

from __future__ import annotations

import io

from .constants import BlockType
from .models import _ParsedBlock
from .tables import _format_kv_fallback, _format_table_semantic


def _parse_pptx(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Extract PPTX title/subtitle/textbox/table/notes in visual order."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ValueError(f"Failed to open PPTX: {exc}") from exc

    blocks: list[_ParsedBlock] = []
    for idx, slide in enumerate(prs.slides, start=1):
        blocks.extend(_parse_pptx_slide(slide, page_number=idx))
    return blocks, len(prs.slides)


def _parse_pptx_slide(slide: object, *, page_number: int) -> list[_ParsedBlock]:
    """Parse a single PPTX slide into ordered blocks."""
    section = f"Slide {page_number}"
    blocks: list[_ParsedBlock] = []
    shapes = sorted(
        slide.shapes,  # type: ignore[attr-defined]
        key=lambda s: (int(getattr(s, "top", 0) or 0), int(getattr(s, "left", 0) or 0)),
    )
    for shape in shapes:
        for b in _pptx_shape_blocks(shape, page_number=page_number, section=section):
            if b.block_type in {"title", "heading"} and b.text.strip():
                section = b.text.strip()
            blocks.append(
                _ParsedBlock(
                    text=b.text,
                    page_number=b.page_number,
                    section=section,
                    heading_level=b.heading_level,
                    block_type=b.block_type,
                )
            )

    notes_text = _pptx_notes_text(slide)
    if notes_text:
        blocks.append(
            _ParsedBlock(
                text=notes_text,
                page_number=page_number,
                section=section,
                block_type="notes",
            )
        )
    return blocks


def _pptx_shape_blocks(
    shape: object,
    *,
    page_number: int,
    section: str,
) -> list[_ParsedBlock]:
    """Parse one PPTX shape into zero or more blocks."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[_ParsedBlock] = []
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.TABLE or getattr(shape, "has_table", False):
        try:
            table = shape.table  # type: ignore[attr-defined]
        except Exception:
            table = None
        if table is not None:
            text = _pptx_table_to_text(table)
            if text.strip():
                out.append(
                    _ParsedBlock(
                        text=text,
                        page_number=page_number,
                        section=section,
                        block_type="table",
                    )
                )
            return out

    if not getattr(shape, "has_text_frame", False):
        return out

    text = (getattr(shape, "text", None) or "").strip()
    if not text:
        return out

    btype: BlockType = "paragraph"
    level: int | None = None
    if getattr(shape, "is_placeholder", False):
        try:
            ph_idx = shape.placeholder_format.idx  # type: ignore[attr-defined]
            ph_type = shape.placeholder_format.type  # type: ignore[attr-defined]
        except Exception:
            ph_idx, ph_type = None, None
        name = str(ph_type).lower() if ph_type is not None else ""
        if ph_idx == 0 or "title" in name and "sub" not in name:
            btype = "title"
            level = 1
        elif "sub" in name:
            btype = "subtitle"
            level = 2

    out.append(
        _ParsedBlock(
            text=text,
            page_number=page_number,
            section=section if btype == "paragraph" else text,
            heading_level=level,
            block_type=btype,
        )
    )
    return out


def _pptx_table_to_text(table: object) -> str:
    rows_data: list[list[str]] = []
    for row in getattr(table, "rows", []) or []:
        cells = [(getattr(c, "text", None) or "").strip() for c in getattr(row, "cells", []) or []]
        if any(cells):
            rows_data.append(cells)
    if not rows_data:
        return ""
    if len(rows_data) >= 2:
        return _format_table_semantic(rows_data[0], rows_data[1:])
    return _format_kv_fallback(rows_data)


def _pptx_notes_text(slide: object) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        notes_slide = slide.notes_slide  # type: ignore[attr-defined]
        text = (notes_slide.notes_text_frame.text or "").strip()
        return text
    except Exception:
        return ""
