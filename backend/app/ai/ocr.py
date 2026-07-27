# =============================================================================
# File: ocr.py
# Module/Service: Pipeline Worker — Document Parsing & Cleaning ([AI])
# Layer: Service
# Purpose: Multi-format document parse + cleaning into layout-aware segments
#          for Chunking / RAG (FR2 Step 3). Not image OCR; not LLM-based.
# Responsibilities:
#   - PDF/DOCX/XLSX/PPTX/TXT → OcrSegment with layout metadata
#   - Heading / table / list detection; paragraph reconstruction
#   - Fuzzy header/footer stripping; Unicode + punctuation normalization
# Dependencies:
#   - PyMuPDF, python-docx, openpyxl, python-pptx (stdlib cleaning only)
# Public Exports:
#   - OcrSegment, CleanedPage, OcrResult, run_ocr_cleaning, EmptyOcrError
# Database/Table: N/A (page_count updated by stage via document_versions)
# Related Modules: app.workers.stages.ocr_cleaning, app.ai.chunking
# Important Notes: No image OCR / no LLM; empty text layer → EmptyOcrError.
# =============================================================================

from __future__ import annotations

import io
import re
import unicodedata
from collections import Counter
from dataclasses import dataclass
from difflib import SequenceMatcher
from typing import Iterator, Literal

from app.models.enums import FileType

# ---------------------------------------------------------------------------
# Constants (no magic numbers in logic)
# ---------------------------------------------------------------------------

BlockType = Literal[
    "paragraph",
    "heading",
    "table",
    "list",
    "caption",
    "title",
    "subtitle",
    "notes",
]

HEADER_FOOTER_MIN_PAGES = 3
HEADER_FOOTER_THRESHOLD = 0.6
HEADER_FOOTER_MIN_LEN = 4
HEADER_FOOTER_MAX_LEN = 120
HEADER_FOOTER_FUZZY_RATIO = 0.88
HEADER_FOOTER_EDGE_LINES = 2

SOFT_BREAK_MAX_GAP_RATIO = 1.35
HARD_BREAK_MIN_GAP_RATIO = 1.75
HEADING_SIZE_RATIO = 1.15
HEADING_MAX_CHARS = 200
TABLE_MIN_ROWS = 2
TABLE_MIN_COLS = 2
SPAN_SPACE_GAP_PT = 1.5
XLSX_MAX_HEADER_SCAN = 20

ZERO_WIDTH_CHARS = (
    "\u200b"  # zero-width space
    "\u200c"  # zero-width non-joiner
    "\u200d"  # zero-width joiner
    "\ufeff"  # BOM / zero-width no-break
    "\u2060"  # word joiner
)

QUOTE_TRANSLATION = str.maketrans(
    {
        "\u2018": "'",
        "\u2019": "'",
        "\u201a": "'",
        "\u201b": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u201e": '"',
        "\u00ab": '"',
        "\u00bb": '"',
    }
)

DASH_TRANSLATION = str.maketrans(
    {
        "\u2010": "-",
        "\u2011": "-",
        "\u2012": "-",
        "\u2013": "-",
        "\u2014": "-",
        "\u2015": "-",
        "\u2212": "-",
    }
)

BULLET_CHARS = frozenset("•●○◆▪▸►‣·∙■□–-*")

# ---------------------------------------------------------------------------
# Precompiled regexes
# ---------------------------------------------------------------------------

_WS_RE = re.compile(r"[ \t]+")
_BLANK_RE = re.compile(r"\n{3,}")
_CTRL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")
_MULTI_SPACE_RE = re.compile(r" {2,}")
_PAGE_NUM_RE = re.compile(
    r"^(?:page|trang|p\.?)\s*\d+(?:\s*(?:of|/|trên)\s*\d+)?$",
    re.IGNORECASE,
)
_STANDALONE_PAGE_RE = re.compile(r"^\d{1,4}$")
_DATE_RE = re.compile(
    r"\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}"
    r"|\d{4}[/-]\d{1,2}[/-]\d{1,2}"
    r"|(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4})\b",
    re.IGNORECASE,
)
_CONFIDENTIAL_RE = re.compile(
    r"\b(?:confidential|internal\s+use\s+only|proprietary|draft)\b",
    re.IGNORECASE,
)
_NUMBERED_HEADING_RE = re.compile(
    r"^(?:"
    r"(?:chapter|section|part|article|appendix)\s+[\dIVXLCDM]+"
    r"|[\dIVXLCDM]+(?:\.\d+){0,4}\.?"
    r")(?:\s+\S.+)?$",
    re.IGNORECASE,
)
_ALL_CAPS_WORD_RE = re.compile(r"^[A-Z0-9][A-Z0-9\s\-:,.&/()]{2,}$")
_LIST_ITEM_RE = re.compile(
    r"^(?:"
    r"[•●○◆▪▸►‣·∙■□]"
    r"|[-*–—]"
    r"|\(?\d{1,3}[.)]"
    r"|[a-zA-Z][.)]"
    r")\s+\S"
)
_CAPTION_RE = re.compile(
    r"^(?:figure|fig\.|table|tbl\.|biểu\s+đồ|bảng)\s*\d+",
    re.IGNORECASE,
)


# ---------------------------------------------------------------------------
# Public exceptions & dataclasses
# ---------------------------------------------------------------------------


class EmptyOcrError(ValueError):
    """Raised when a file yields no extractable text after cleaning.

    Typical cause: scanned PDF without a text layer (image OCR is out of scope
    for this stage).
    """


@dataclass(frozen=True, slots=True)
class OcrSegment:
    """Normalized text unit shared by all input formats.

    Core fields are required by existing consumers. Layout metadata fields are
    optional and backward-compatible (default ``None``).
    """

    text: str
    order_index: int
    page_number: int | None = None
    section: str | None = None
    heading_level: int | None = None
    block_type: BlockType | None = None
    bbox: tuple[float, float, float, float] | None = None
    language: str | None = None
    font_size: float | None = None
    font_name: str | None = None
    is_bold: bool | None = None


@dataclass(frozen=True, slots=True)
class CleanedPage:
    """Legacy page view for structure-aware chunking (derived from segments)."""

    page_number: int
    text: str
    section: str | None = None


@dataclass(frozen=True, slots=True)
class OcrMetrics:
    """Pipeline observability counters (not persisted unless a stage chooses to)."""

    page_count: int
    char_count: int
    segment_count: int
    heading_count: int
    table_count: int


@dataclass(frozen=True, slots=True)
class OcrResult:
    """OCR/cleaning output for one document version."""

    segments: list[OcrSegment]
    page_count: int
    char_count: int

    @property
    def pages(self) -> list[CleanedPage]:
        """Adapt segments to ``CleanedPage`` for the chunking module."""
        return [
            CleanedPage(
                page_number=(
                    seg.page_number if seg.page_number is not None else seg.order_index + 1
                ),
                text=seg.text,
                section=seg.section,
            )
            for seg in self.segments
        ]

    @property
    def segment_count(self) -> int:
        return len(self.segments)

    @property
    def heading_count(self) -> int:
        return sum(
            1
            for s in self.segments
            if s.block_type == "heading" or (s.heading_level is not None and s.heading_level > 0)
        )

    @property
    def table_count(self) -> int:
        return sum(1 for s in self.segments if s.block_type == "table")

    @property
    def metrics(self) -> OcrMetrics:
        return OcrMetrics(
            page_count=self.page_count,
            char_count=self.char_count,
            segment_count=self.segment_count,
            heading_count=self.heading_count,
            table_count=self.table_count,
        )


# ---------------------------------------------------------------------------
# Internal parsed block (stateless intermediate)
# ---------------------------------------------------------------------------


@dataclass(frozen=True, slots=True)
class _ParsedBlock:
    """Intermediate parse unit before final cleaning / indexing."""

    text: str
    page_number: int | None = None
    section: str | None = None
    heading_level: int | None = None
    block_type: BlockType = "paragraph"
    bbox: tuple[float, float, float, float] | None = None
    font_size: float | None = None
    font_name: str | None = None
    is_bold: bool | None = None


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def run_ocr_cleaning(*, file_type: FileType, data: bytes) -> OcrResult:
    """Parse and clean a document into normalized OCR segments.

    Args:
        file_type: Declared document type from ``documents.file_type``.
        data: Raw file bytes from object storage.

    Returns:
        ``OcrResult`` with segments, physical/logical ``page_count``, and
        aggregate ``char_count``. Use ``.metrics`` for pipeline observability.

    Raises:
        EmptyOcrError: No non-empty text after cleaning.
        ValueError: Unsupported ``file_type``.
        Exception: Propagated from format parsers (corrupt files).
    """
    if file_type == FileType.pdf:
        blocks, page_count = _parse_pdf(data)
    elif file_type == FileType.docx:
        blocks, page_count = _parse_docx(data)
    elif file_type == FileType.xlsx:
        blocks, page_count = _parse_xlsx(data)
    elif file_type == FileType.pptx:
        blocks, page_count = _parse_pptx(data)
    elif file_type == FileType.txt:
        blocks, page_count = _parse_txt(data)
    else:
        raise ValueError(f"Unsupported file_type: {file_type}")

    stripped = _strip_repeated_headers_footers(blocks)
    segments = list(_blocks_to_segments(stripped))

    char_count = sum(len(s.text) for s in segments)
    if not segments or char_count == 0:
        raise EmptyOcrError(
            "No extractable text after OCR/cleaning. If this is a scanned PDF, "
            "it has no text layer — image OCR is not enabled in this stage. "
            "Re-upload a text-based document or a PDF with an embedded text layer."
        )

    return OcrResult(segments=segments, page_count=page_count, char_count=char_count)


def _blocks_to_segments(blocks: list[_ParsedBlock]) -> Iterator[OcrSegment]:
    """Clean blocks and assign stable ``order_index`` values."""
    order = 0
    for block in blocks:
        cleaned = _clean_text(block.text)
        if not cleaned:
            continue
        # Tables / headings already emitted as semantic units — keep structure.
        if block.block_type in {"table", "heading", "title", "subtitle", "notes", "caption"}:
            yield OcrSegment(
                text=cleaned,
                order_index=order,
                page_number=block.page_number,
                section=block.section,
                heading_level=block.heading_level,
                block_type=block.block_type,
                bbox=block.bbox,
                font_size=block.font_size,
                font_name=block.font_name,
                is_bold=block.is_bold,
            )
            order += 1
            continue

        # Soft wraps → spaces; blank lines remain hard paragraph breaks.
        cleaned = _join_soft_lines(cleaned.split("\n"))
        for para in _split_paragraphs(cleaned):
            yield OcrSegment(
                text=para,
                order_index=order,
                page_number=block.page_number,
                section=block.section,
                heading_level=block.heading_level,
                block_type=block.block_type or _infer_block_type(para),
                bbox=block.bbox,
                font_size=block.font_size,
                font_name=block.font_name,
                is_bold=block.is_bold,
            )
            order += 1


# ---------------------------------------------------------------------------
# Text cleaning
# ---------------------------------------------------------------------------


def _clean_text(text: str) -> str:
    """Normalize encoding, punctuation, and whitespace for RAG-friendly text."""
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    for ch in ZERO_WIDTH_CHARS:
        if ch in text:
            text = text.replace(ch, "")
    text = text.replace("\u00a0", " ")
    text = text.translate(QUOTE_TRANSLATION)
    text = text.translate(DASH_TRANSLATION)
    text = _normalize_bullets(text)
    text = _CTRL_RE.sub("", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _WS_RE.sub(" ", text)
    text = _BLANK_RE.sub("\n\n", text)
    return text.strip()


def _normalize_bullets(text: str) -> str:
    """Map decorative bullet glyphs to a single ASCII bullet."""
    if not text:
        return text
    out: list[str] = []
    for line in text.split("\n"):
        stripped = line.lstrip()
        if stripped and stripped[0] in BULLET_CHARS:
            indent = line[: len(line) - len(stripped)]
            rest = stripped[1:].lstrip()
            out.append(f"{indent}* {rest}" if rest else f"{indent}*")
        else:
            out.append(line)
    return "\n".join(out)


def _split_paragraphs(text: str) -> list[str]:
    """Split a block into paragraph-sized pieces (hard breaks only)."""
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        return []
    return [p.strip() for p in re.split(r"\n\s*\n+", text) if p.strip()]


def _infer_block_type(text: str) -> BlockType:
    if _CAPTION_RE.match(text):
        return "caption"
    if _LIST_ITEM_RE.match(text):
        return "list"
    if _looks_like_heading_text(text):
        return "heading"
    return "paragraph"


# ---------------------------------------------------------------------------
# Heading detection
# ---------------------------------------------------------------------------


def _looks_like_heading_text(text: str) -> bool:
    """Heuristic heading detection from plain text (no font metadata)."""
    candidate = text.strip()
    if not candidate or len(candidate) > HEADING_MAX_CHARS:
        return False
    if "\n" in candidate:
        return False
    if _NUMBERED_HEADING_RE.match(candidate):
        return True
    letters = [c for c in candidate if c.isalpha()]
    if letters and _ALL_CAPS_WORD_RE.match(candidate) and len(letters) >= 3:
        return True
    return False


def _heading_level_from_text(text: str) -> int:
    """Infer heading depth from numbered prefixes (1 → 1, 1.1 → 2)."""
    match = re.match(r"^(\d+(?:\.\d+)*)", text.strip())
    if match:
        return min(match.group(1).count(".") + 1, 6)
    upper = text.strip().upper()
    if upper.startswith(("CHAPTER", "PART", "APPENDIX")):
        return 1
    if upper.startswith("SECTION"):
        return 2
    return 1


def _detect_heading_from_font(
    text: str,
    *,
    font_size: float | None,
    is_bold: bool | None,
    median_size: float,
) -> tuple[bool, int | None]:
    """Detect heading using font size / weight plus text heuristics."""
    candidate = text.strip()
    if not candidate or len(candidate) > HEADING_MAX_CHARS:
        return False, None

    size_boost = False
    if font_size is not None and median_size > 0:
        size_boost = font_size >= median_size * HEADING_SIZE_RATIO

    text_hit = _looks_like_heading_text(candidate)
    bold_short = bool(is_bold) and len(candidate) <= 120 and not candidate.endswith(".")

    if size_boost or text_hit or bold_short:
        level = _heading_level_from_text(candidate)
        if size_boost and font_size is not None and median_size > 0:
            ratio = font_size / median_size
            if ratio >= 1.6:
                level = 1
            elif ratio >= 1.35:
                level = min(level, 2)
        return True, level
    return False, None


# ---------------------------------------------------------------------------
# Paragraph reconstruction (soft vs hard breaks)
# ---------------------------------------------------------------------------


def _is_soft_line_break(prev: str, nxt: str) -> bool:
    """True when two visual lines should merge into one paragraph."""
    prev = prev.rstrip()
    nxt = nxt.lstrip()
    if not prev or not nxt:
        return False
    if prev.endswith(("-", "\u2010", "\u2011")):
        return True
    if prev[-1] in ".!?:;":
        return False
    if nxt[0].islower():
        return True
    if prev[-1] in ",;:" or prev[-1].isalnum():
        return True
    return False


def _join_soft_lines(lines: list[str]) -> str:
    """Merge soft wraps; preserve hard paragraph gaps as blank lines."""
    if not lines:
        return ""
    parts: list[str] = []
    buf = lines[0].rstrip()
    for nxt in lines[1:]:
        nxt_s = nxt.strip()
        if not nxt_s:
            if buf:
                parts.append(buf)
                buf = ""
            continue
        if not buf:
            buf = nxt_s
            continue
        if _is_soft_line_break(buf, nxt_s):
            if buf.endswith(("-", "\u2010", "\u2011")) and nxt_s and nxt_s[0].islower():
                buf = buf[:-1] + nxt_s
            else:
                buf = f"{buf} {nxt_s}"
        else:
            parts.append(buf)
            buf = nxt_s
    if buf:
        parts.append(buf)
    return "\n\n".join(parts)


def _reconstruct_paragraphs_from_lines(
    line_texts: list[str],
    line_sizes: list[float],
    line_bboxes: list[tuple[float, float, float, float]],
) -> list[str]:
    """Rebuild paragraphs using vertical gaps and soft-wrap heuristics."""
    if not line_texts:
        return []

    groups: list[list[str]] = [[line_texts[0]]]
    for i in range(1, len(line_texts)):
        prev_bbox = line_bboxes[i - 1]
        cur_bbox = line_bboxes[i]
        avg_size = (line_sizes[i - 1] + line_sizes[i]) / 2.0 or 12.0
        gap = cur_bbox[1] - prev_bbox[3]
        prev_text = line_texts[i - 1]
        cur_text = line_texts[i]

        hard = gap > avg_size * HARD_BREAK_MIN_GAP_RATIO
        soft = (
            gap <= avg_size * SOFT_BREAK_MAX_GAP_RATIO
            and _is_soft_line_break(prev_text, cur_text)
        )
        if hard or not soft:
            # Sentence end + capital start with moderate gap → new paragraph
            groups.append([cur_text])
        else:
            groups[-1].append(cur_text)

    return [_join_soft_lines(g) for g in groups if any(t.strip() for t in g)]


# ---------------------------------------------------------------------------
# Table formatting helpers
# ---------------------------------------------------------------------------


def _format_table_semantic(
    headers: list[str],
    rows: list[list[str]],
    *,
    title: str | None = None,
) -> str:
    """Render a table as ``Column : Value`` rows for better embeddings."""
    parts: list[str] = []
    if title:
        parts.append(title)
    headers = [h.strip() for h in headers]
    for idx, row in enumerate(rows, start=1):
        parts.append(f"Row {idx}")
        for col_i, cell in enumerate(row):
            cell = (cell or "").strip()
            if not cell:
                continue
            if col_i < len(headers) and headers[col_i]:
                parts.append(f"{headers[col_i]} = {cell}")
            else:
                parts.append(f"Column{col_i + 1} = {cell}")
    return "\n".join(parts)


def _format_kv_fallback(rows: list[list[str]]) -> str:
    """Fallback when headers are unknown: ``ColumnA : ValueA`` per cell pair."""
    parts: list[str] = []
    for idx, row in enumerate(rows, start=1):
        cells = [c.strip() for c in row if c and c.strip()]
        if not cells:
            continue
        parts.append(f"Row {idx}")
        if len(cells) == 2:
            parts.append(f"{cells[0]} : {cells[1]}")
        else:
            for i, cell in enumerate(cells, start=1):
                parts.append(f"Column{i} : {cell}")
    return "\n".join(parts)


def _detect_table_from_aligned_lines(
    lines: list[tuple[str, list[tuple[float, str]]]],
) -> str | None:
    """Detect column-aligned text lines and format as semantic table text.

    Args:
        lines: Each item is ``(full_line_text, [(x0, span_text), ...])``.
    """
    if len(lines) < TABLE_MIN_ROWS:
        return None

    col_counts = [len(spans) for _, spans in lines if len(spans) >= TABLE_MIN_COLS]
    if len(col_counts) < TABLE_MIN_ROWS:
        return None

    # Dominant column count
    n_cols, freq = Counter(col_counts).most_common(1)[0]
    if n_cols < TABLE_MIN_COLS or freq < TABLE_MIN_ROWS:
        return None

    matrix: list[list[str]] = []
    for _, spans in lines:
        if len(spans) != n_cols:
            continue
        matrix.append([t.strip() for _, t in spans])
    if len(matrix) < TABLE_MIN_ROWS:
        return None

    headers = matrix[0]
    body = matrix[1:]
    if body and all(headers):
        return _format_table_semantic(headers, body)
    return _format_kv_fallback(matrix)


# ---------------------------------------------------------------------------
# Fuzzy header / footer stripping
# ---------------------------------------------------------------------------


def _normalize_header_footer_key(text: str) -> str:
    """Normalize a line for fuzzy header/footer comparison."""
    key = text.strip().lower()
    key = _PAGE_NUM_RE.sub("page #", key)
    key = _DATE_RE.sub("DATE", key)
    key = _CONFIDENTIAL_RE.sub("CONFIDENTIAL", key)
    key = re.sub(r"\d+", "#", key)
    key = _MULTI_SPACE_RE.sub(" ", key).strip()
    return key


def _is_boilerplate_line(text: str) -> bool:
    """Likely page chrome (page numbers, confidentiality, short dates)."""
    t = text.strip()
    if not t:
        return False
    if _PAGE_NUM_RE.match(t) or _STANDALONE_PAGE_RE.match(t):
        return True
    if _CONFIDENTIAL_RE.search(t) and len(t) <= 80:
        return True
    if _DATE_RE.fullmatch(t):
        return True
    return False


def _fuzzy_match(a: str, b: str, *, threshold: float = HEADER_FOOTER_FUZZY_RATIO) -> bool:
    if a == b:
        return True
    na, nb = _normalize_header_footer_key(a), _normalize_header_footer_key(b)
    if not na or not nb:
        return False
    if na == nb:
        return True
    return SequenceMatcher(None, na, nb).ratio() >= threshold


def _majority_fuzzy_line(candidates: list[str], *, threshold: float) -> str | None:
    """Return a representative line present (fuzzily) on ≥ threshold of pages."""
    usable = [
        c.strip()
        for c in candidates
        if HEADER_FOOTER_MIN_LEN <= len(c.strip()) <= HEADER_FOOTER_MAX_LEN
        or _is_boilerplate_line(c)
    ]
    if not usable:
        return None

    clusters: list[list[str]] = []
    for line in usable:
        placed = False
        for cluster in clusters:
            if _fuzzy_match(cluster[0], line):
                cluster.append(line)
                placed = True
                break
        if not placed:
            clusters.append([line])

    best = max(clusters, key=len)
    # Denominator = number of candidate slots (one per page that had a line)
    if len(best) / max(len(candidates), 1) >= threshold:
        # Prefer the most common exact form inside the cluster
        return Counter(best).most_common(1)[0][0]
    return None


def _line_matches_drop(line: str, drop: str | None) -> bool:
    if drop is None:
        return False
    candidate = line.strip()
    if not candidate:
        return False
    if _is_boilerplate_line(candidate) and _is_boilerplate_line(drop):
        return _fuzzy_match(candidate, drop, threshold=0.75)
    return _fuzzy_match(candidate, drop)


def _strip_repeated_headers_footers(blocks: list[_ParsedBlock]) -> list[_ParsedBlock]:
    """Remove fuzzy-repeated header/footer lines across a majority of pages."""
    page_lines: dict[int, list[str]] = {}
    page_order: list[int] = []
    for block in blocks:
        if block.page_number is None:
            continue
        if block.page_number not in page_lines:
            page_order.append(block.page_number)
            page_lines[block.page_number] = []
        for ln in block.text.splitlines():
            if ln.strip():
                page_lines[block.page_number].append(ln.strip())

    if len(page_order) < HEADER_FOOTER_MIN_PAGES:
        return blocks

    headers: list[str] = []
    footers: list[str] = []
    for pn in page_order:
        lines = page_lines.get(pn) or []
        if not lines:
            continue
        headers.append(lines[0])
        footers.append(lines[-1])

    drop_header = _majority_fuzzy_line(headers, threshold=HEADER_FOOTER_THRESHOLD)
    drop_footer = _majority_fuzzy_line(footers, threshold=HEADER_FOOTER_THRESHOLD)
    if drop_header is None and drop_footer is None:
        # Still strip obvious per-page boilerplate at edges
        drop_header = drop_footer = None

    cleaned: list[_ParsedBlock] = []
    for block in blocks:
        if block.page_number is None or block.block_type in {"table", "heading", "title"}:
            cleaned.append(block)
            continue
        lines = block.text.splitlines()
        if not lines:
            cleaned.append(block)
            continue
        kept: list[str] = []
        n = len(lines)
        for i, ln in enumerate(lines):
            at_top = i < HEADER_FOOTER_EDGE_LINES
            at_bottom = i >= n - HEADER_FOOTER_EDGE_LINES
            if at_top and _line_matches_drop(ln, drop_header):
                continue
            if at_bottom and _line_matches_drop(ln, drop_footer):
                continue
            if (at_top or at_bottom) and _is_boilerplate_line(ln) and len(page_order) >= HEADER_FOOTER_MIN_PAGES:
                # Page numbers / confidential markers on edges
                continue
            kept.append(ln)
        text = "\n".join(kept)
        if text.strip() == block.text.strip():
            cleaned.append(block)
        else:
            cleaned.append(
                _ParsedBlock(
                    text=text,
                    page_number=block.page_number,
                    section=block.section,
                    heading_level=block.heading_level,
                    block_type=block.block_type,
                    bbox=block.bbox,
                    font_size=block.font_size,
                    font_name=block.font_name,
                    is_bold=block.is_bold,
                )
            )
    return cleaned


# ---------------------------------------------------------------------------
# PDF parser (layout-aware via get_text("dict"))
# ---------------------------------------------------------------------------


def _parse_pdf(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Extract PDF text with block/line/span layout; rebuild reading order."""
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise ValueError(f"Failed to open PDF: {exc}") from exc

    try:
        if len(doc) == 0:
            return [], 0

        # First pass: collect font sizes for median body size
        all_sizes: list[float] = []
        page_dicts: list[dict] = []
        for page in doc:
            try:
                page_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            except Exception as exc:
                raise ValueError(f"Failed to extract PDF page layout: {exc}") from exc
            page_dicts.append(page_dict)
            for block in page_dict.get("blocks") or []:
                if block.get("type", 0) != 0:
                    continue
                for line in block.get("lines") or []:
                    for span in line.get("spans") or []:
                        size = float(span.get("size") or 0)
                        if size > 0:
                            all_sizes.append(size)

        median_size = _median(all_sizes) if all_sizes else 12.0
        blocks_out: list[_ParsedBlock] = []
        current_section: str | None = None

        for page_no, page_dict in enumerate(page_dicts, start=1):
            page_blocks = _pdf_page_to_blocks(
                page_dict,
                page_number=page_no,
                median_size=median_size,
                current_section=current_section,
            )
            for b in page_blocks:
                blocks_out.append(b)
                if b.block_type == "heading" and b.text.strip():
                    current_section = b.text.strip()

        return blocks_out, len(doc)
    finally:
        doc.close()


@dataclass(frozen=True, slots=True)
class _PdfLineBundle:
    """Extracted lines from one PDF text block."""

    texts: list[str]
    sizes: list[float]
    bboxes: list[tuple[float, float, float, float]]
    span_cols: list[tuple[str, list[tuple[float, str]]]]
    block_bbox: tuple[float, float, float, float]
    font_size: float
    font_name: str
    is_bold: bool


def _pdf_page_to_blocks(
    page_dict: dict,
    *,
    page_number: int,
    median_size: float,
    current_section: str | None,
) -> list[_ParsedBlock]:
    """Convert one PDF page dict into ordered parsed blocks."""
    result: list[_ParsedBlock] = []
    section = current_section

    for block in page_dict.get("blocks") or []:
        if block.get("type", 0) != 0:
            continue
        bundle = _extract_pdf_line_bundle(block)
        if bundle is None:
            continue
        emitted, section = _emit_pdf_block(
            bundle,
            page_number=page_number,
            median_size=median_size,
            section=section,
        )
        result.extend(emitted)
    return result


def _extract_pdf_line_bundle(block: dict) -> _PdfLineBundle | None:
    """Pull line texts / fonts / columns from a PyMuPDF text block."""
    lines_raw = block.get("lines") or []
    if not lines_raw:
        return None

    texts: list[str] = []
    sizes: list[float] = []
    bboxes: list[tuple[float, float, float, float]] = []
    span_cols: list[tuple[str, list[tuple[float, str]]]] = []
    block_bold = False
    block_font = ""
    block_size = 0.0

    for line in lines_raw:
        text, size, font, bold, spans_x = _merge_pdf_spans(line.get("spans") or [])
        if not text.strip():
            continue
        bbox = tuple(float(x) for x in (line.get("bbox") or (0, 0, 0, 0)))
        texts.append(text)
        sizes.append(size)
        bboxes.append(bbox)  # type: ignore[arg-type]
        span_cols.append((text, spans_x))
        if size >= block_size:
            block_size = size
            block_font = font
            block_bold = bold

    if not texts:
        return None

    block_bbox = tuple(float(x) for x in (block.get("bbox") or (0, 0, 0, 0)))
    return _PdfLineBundle(
        texts=texts,
        sizes=sizes,
        bboxes=bboxes,
        span_cols=span_cols,
        block_bbox=block_bbox,  # type: ignore[arg-type]
        font_size=block_size,
        font_name=block_font,
        is_bold=block_bold,
    )


def _emit_pdf_block(
    bundle: _PdfLineBundle,
    *,
    page_number: int,
    median_size: float,
    section: str | None,
) -> tuple[list[_ParsedBlock], str | None]:
    """Emit table / heading / paragraph blocks from one PDF line bundle."""
    table_text = _detect_table_from_aligned_lines(bundle.span_cols)
    if table_text:
        return (
            [
                _ParsedBlock(
                    text=table_text,
                    page_number=page_number,
                    section=section,
                    block_type="table",
                    bbox=bundle.block_bbox,
                    font_size=bundle.font_size or None,
                    font_name=bundle.font_name or None,
                    is_bold=bundle.is_bold,
                )
            ],
            section,
        )

    probe = bundle.texts[0] if len(bundle.texts) != 1 else bundle.texts[0]
    is_heading, level = _detect_heading_from_font(
        probe,
        font_size=bundle.font_size,
        is_bold=bundle.is_bold,
        median_size=median_size,
    )
    if is_heading and (len(bundle.texts) == 1 or _looks_like_heading_text(bundle.texts[0])):
        heading_text = bundle.texts[0].strip()
        section = heading_text
        out: list[_ParsedBlock] = [
            _ParsedBlock(
                text=heading_text,
                page_number=page_number,
                section=section,
                heading_level=level,
                block_type="heading",
                bbox=bundle.bboxes[0],
                font_size=bundle.font_size or None,
                font_name=bundle.font_name or None,
                is_bold=bundle.is_bold,
            )
        ]
        if len(bundle.texts) > 1:
            for para in _reconstruct_paragraphs_from_lines(
                bundle.texts[1:], bundle.sizes[1:], bundle.bboxes[1:]
            ):
                out.append(
                    _ParsedBlock(
                        text=para,
                        page_number=page_number,
                        section=section,
                        block_type="paragraph",
                        bbox=bundle.block_bbox,
                        font_size=bundle.font_size or None,
                        font_name=bundle.font_name or None,
                        is_bold=False,
                    )
                )
        return out, section

    out = []
    for para in _reconstruct_paragraphs_from_lines(bundle.texts, bundle.sizes, bundle.bboxes):
        btype: BlockType = "list" if _LIST_ITEM_RE.match(para.strip()) else "paragraph"
        out.append(
            _ParsedBlock(
                text=para,
                page_number=page_number,
                section=section,
                block_type=btype,
                bbox=bundle.block_bbox,
                font_size=bundle.font_size or None,
                font_name=bundle.font_name or None,
                is_bold=bundle.is_bold,
            )
        )
    return out, section


def _merge_pdf_spans(
    spans: list[dict],
) -> tuple[str, float, str, bool, list[tuple[float, str]]]:
    """Merge spans on one line; return text, dominant font meta, and x-columns."""
    if not spans:
        return "", 0.0, "", False, []

    parts: list[str] = []
    cols: list[tuple[float, str]] = []
    max_size = 0.0
    font_name = ""
    is_bold = False
    prev_x1: float | None = None

    for span in spans:
        raw = span.get("text") or ""
        if not raw:
            continue
        bbox = span.get("bbox") or (0.0, 0.0, 0.0, 0.0)
        x0, _, x1, _ = (float(bbox[0]), float(bbox[1]), float(bbox[2]), float(bbox[3]))
        size = float(span.get("size") or 0.0)
        font = str(span.get("font") or "")
        flags = int(span.get("flags") or 0)
        bold = bool(flags & 2**4) or ("bold" in font.lower())

        if size >= max_size:
            max_size = size
            font_name = font
            is_bold = bold

        if parts and prev_x1 is not None:
            gap = x0 - prev_x1
            if gap > SPAN_SPACE_GAP_PT and not parts[-1].endswith(" ") and not raw.startswith(" "):
                parts.append(" ")
        parts.append(raw)
        prev_x1 = x1
        cols.append((x0, raw.strip()))

    return "".join(parts), max_size, font_name, is_bold, cols


def _median(values: list[float]) -> float:
    if not values:
        return 0.0
    ordered = sorted(values)
    mid = len(ordered) // 2
    if len(ordered) % 2:
        return ordered[mid]
    return (ordered[mid - 1] + ordered[mid]) / 2.0


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------


def _parse_docx(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Extract DOCX body in order: headings, paragraphs, lists, tables, captions."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph

    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise ValueError(f"Failed to open DOCX: {exc}") from exc

    blocks: list[_ParsedBlock] = []
    current_section: str | None = None
    logical_page = 1

    for child in document.element.body:
        tag = child.tag
        if tag == qn("w:p"):
            para = Paragraph(child, document)
            parsed = _docx_paragraph_block(para, logical_page, current_section)
            if parsed is None:
                continue
            if parsed.block_type == "heading":
                current_section = parsed.text.strip()
                logical_page += 1
            blocks.append(
                _ParsedBlock(
                    text=parsed.text,
                    page_number=logical_page,
                    section=current_section if parsed.block_type != "heading" else parsed.text.strip(),
                    heading_level=parsed.heading_level,
                    block_type=parsed.block_type,
                    is_bold=parsed.is_bold,
                )
            )
            if parsed.block_type == "heading":
                current_section = parsed.text.strip()
        elif tag == qn("w:tbl"):
            table = DocxTable(child, document)
            text = _docx_table_to_text(table)
            if text.strip():
                blocks.append(
                    _ParsedBlock(
                        text=text,
                        page_number=logical_page,
                        section=current_section,
                        block_type="table",
                    )
                )

    page_count = max((b.page_number or 1 for b in blocks), default=0)
    return blocks, page_count


def _docx_paragraph_block(
    para: object,
    logical_page: int,
    current_section: str | None,
) -> _ParsedBlock | None:
    """Map one DOCX paragraph to a parsed block (heading/list/caption/body)."""
    text = getattr(para, "text", None) or ""
    if not text.strip():
        return None

    style = getattr(para, "style", None)
    style_name = ((style.name or "") if style else "").lower()
    is_bold = _docx_paragraph_is_bold(para)

    if "heading" in style_name or style_name.startswith("title"):
        level = _docx_heading_level(style_name)
        btype: BlockType = "title" if style_name == "title" else "heading"
        return _ParsedBlock(
            text=text.strip(),
            page_number=logical_page,
            section=text.strip(),
            heading_level=level,
            block_type=btype,
            is_bold=True,
        )

    if "caption" in style_name or _CAPTION_RE.match(text.strip()):
        return _ParsedBlock(
            text=text.strip(),
            page_number=logical_page,
            section=current_section,
            block_type="caption",
            is_bold=is_bold,
        )

    if "list" in style_name or _LIST_ITEM_RE.match(text.strip()):
        return _ParsedBlock(
            text=text.strip(),
            page_number=logical_page,
            section=current_section,
            block_type="list",
            is_bold=is_bold,
        )

    return _ParsedBlock(
        text=text.strip(),
        page_number=logical_page,
        section=current_section,
        block_type="paragraph",
        is_bold=is_bold,
    )


def _docx_heading_level(style_name: str) -> int:
    match = re.search(r"(\d+)", style_name)
    if match:
        return min(max(int(match.group(1)), 1), 6)
    if style_name == "title":
        return 1
    return 1


def _docx_paragraph_is_bold(para: object) -> bool | None:
    runs = getattr(para, "runs", None) or []
    if not runs:
        return None
    bold_flags = [bool(getattr(r, "bold", False)) for r in runs if (getattr(r, "text", None) or "").strip()]
    if not bold_flags:
        return None
    return all(bold_flags)


def _docx_table_to_text(table: object) -> str:
    rows_data: list[list[str]] = []
    for row in getattr(table, "rows", []) or []:
        cells = []
        for cell in getattr(row, "cells", []) or []:
            cells.append((getattr(cell, "text", None) or "").strip())
        if any(cells):
            rows_data.append(cells)
    if not rows_data:
        return ""
    if len(rows_data) >= 2:
        return _format_table_semantic(rows_data[0], rows_data[1:])
    return _format_kv_fallback(rows_data)


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


def _parse_pptx(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Extract PPTX title/subtitle/textbox/table/notes in visual order."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ValueError(f"Failed to open PPTX: {exc}") from exc

    blocks: list[_ParsedBlock] = []
    for idx, slide in enumerate(prs.slides, start=1):
        blocks.extend(_parse_pptx_slide(slide, page_number=idx))
    return blocks, len(prs.slides)


def _parse_pptx_slide(slide: object, *, page_number: int) -> list[_ParsedBlock]:
    """Parse a single PPTX slide into ordered blocks."""
    section = f"Slide {page_number}"
    blocks: list[_ParsedBlock] = []
    shapes = sorted(
        slide.shapes,  # type: ignore[attr-defined]
        key=lambda s: (int(getattr(s, "top", 0) or 0), int(getattr(s, "left", 0) or 0)),
    )
    for shape in shapes:
        for b in _pptx_shape_blocks(shape, page_number=page_number, section=section):
            if b.block_type in {"title", "heading"} and b.text.strip():
                section = b.text.strip()
            blocks.append(
                _ParsedBlock(
                    text=b.text,
                    page_number=b.page_number,
                    section=section,
                    heading_level=b.heading_level,
                    block_type=b.block_type,
                )
            )

    notes_text = _pptx_notes_text(slide)
    if notes_text:
        blocks.append(
            _ParsedBlock(
                text=notes_text,
                page_number=page_number,
                section=section,
                block_type="notes",
            )
        )
    return blocks


def _pptx_shape_blocks(
    shape: object,
    *,
    page_number: int,
    section: str,
) -> list[_ParsedBlock]:
    """Parse one PPTX shape into zero or more blocks."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[_ParsedBlock] = []
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.TABLE or getattr(shape, "has_table", False):
        try:
            table = shape.table  # type: ignore[attr-defined]
        except Exception:
            table = None
        if table is not None:
            text = _pptx_table_to_text(table)
            if text.strip():
                out.append(
                    _ParsedBlock(
                        text=text,
                        page_number=page_number,
                        section=section,
                        block_type="table",
                    )
                )
            return out

    if not getattr(shape, "has_text_frame", False):
        return out

    text = (getattr(shape, "text", None) or "").strip()
    if not text:
        return out

    btype: BlockType = "paragraph"
    level: int | None = None
    if getattr(shape, "is_placeholder", False):
        try:
            ph_idx = shape.placeholder_format.idx  # type: ignore[attr-defined]
            ph_type = shape.placeholder_format.type  # type: ignore[attr-defined]
        except Exception:
            ph_idx, ph_type = None, None
        name = str(ph_type).lower() if ph_type is not None else ""
        if ph_idx == 0 or "title" in name and "sub" not in name:
            btype = "title"
            level = 1
        elif "sub" in name:
            btype = "subtitle"
            level = 2

    out.append(
        _ParsedBlock(
            text=text,
            page_number=page_number,
            section=section if btype == "paragraph" else text,
            heading_level=level,
            block_type=btype,
        )
    )
    return out


def _pptx_table_to_text(table: object) -> str:
    rows_data: list[list[str]] = []
    for row in getattr(table, "rows", []) or []:
        cells = [(getattr(c, "text", None) or "").strip() for c in getattr(row, "cells", []) or []]
        if any(cells):
            rows_data.append(cells)
    if not rows_data:
        return ""
    if len(rows_data) >= 2:
        return _format_table_semantic(rows_data[0], rows_data[1:])
    return _format_kv_fallback(rows_data)


def _pptx_notes_text(slide: object) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        notes_slide = slide.notes_slide  # type: ignore[attr-defined]
        text = (notes_slide.notes_text_frame.text or "").strip()
        return text
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# XLSX parser
# ---------------------------------------------------------------------------


def _parse_xlsx(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Extract XLSX sheets as semantic row/column text (not pipe-joined)."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ValueError(f"Failed to open XLSX: {exc}") from exc

    blocks: list[_ParsedBlock] = []
    try:
        for idx, sheet in enumerate(wb.worksheets, start=1):
            rows = _xlsx_collect_rows(sheet)
            if not rows:
                continue
            text = _xlsx_rows_to_semantic(rows, sheet_title=sheet.title)
            if text.strip():
                blocks.append(
                    _ParsedBlock(
                        text=text,
                        page_number=idx,
                        section=sheet.title,
                        block_type="table",
                    )
                )
        return blocks, len(wb.worksheets)
    finally:
        wb.close()


def _xlsx_collect_rows(sheet: object) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in sheet.iter_rows(values_only=True):  # type: ignore[attr-defined]
        cells = ["" if c is None else str(c).strip() for c in row]
        # Trim trailing empties
        while cells and not cells[-1]:
            cells.pop()
        if any(cells):
            rows.append(cells)
    return rows


def _xlsx_rows_to_semantic(rows: list[list[str]], *, sheet_title: str) -> str:
    """Convert sheet rows to ``Row N / Col = Val`` semantic text."""
    if not rows:
        return ""

    header_idx = _xlsx_guess_header_row(rows)
    if header_idx is None:
        return _format_kv_fallback(rows)

    headers = rows[header_idx]
    body = rows[header_idx + 1 :]
    return _format_table_semantic(headers, body, title=f"Sheet: {sheet_title}")


def _xlsx_guess_header_row(rows: list[list[str]]) -> int | None:
    """Pick the first non-empty row with mostly non-numeric labels as header."""
    limit = min(len(rows), XLSX_MAX_HEADER_SCAN)
    for i in range(limit):
        row = rows[i]
        non_empty = [c for c in row if c]
        if len(non_empty) < 1:
            continue
        numericish = sum(1 for c in non_empty if _looks_numeric(c))
        if numericish / len(non_empty) <= 0.4:
            return i
    return 0 if rows else None


def _looks_numeric(value: str) -> bool:
    try:
        float(value.replace(",", "").replace("%", ""))
        return True
    except ValueError:
        return False


# ---------------------------------------------------------------------------
# TXT parser
# ---------------------------------------------------------------------------


def _parse_txt(data: bytes) -> tuple[list[_ParsedBlock], int]:
    """Decode TXT and keep paragraph structure (cleaning applied later)."""
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception as exc:
        raise ValueError(f"Failed to decode TXT: {exc}") from exc

    paragraphs = _split_paragraphs(text)
    if not paragraphs:
        return [_ParsedBlock(text="", page_number=1, block_type="paragraph")], 0

    current_section: str | None = None
    blocks: list[_ParsedBlock] = []
    for para in paragraphs:
        if _looks_like_heading_text(para):
            current_section = para.strip()
            blocks.append(
                _ParsedBlock(
                    text=para,
                    page_number=1,
                    section=current_section,
                    heading_level=_heading_level_from_text(para),
                    block_type="heading",
                )
            )
        else:
            btype: BlockType = "list" if _LIST_ITEM_RE.match(para.strip()) else "paragraph"
            blocks.append(
                _ParsedBlock(
                    text=para,
                    page_number=1,
                    section=current_section,
                    block_type=btype,
                )
            )
    return blocks, 1
